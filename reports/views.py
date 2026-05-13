from django.shortcuts import render, redirect
from django.contrib.auth.decorators import login_required
from pptx import Presentation
from .ai_utils import generate_newsletter_narrative
from .models import MonthlyNewsletter
from django.utils import timezone
from io import BytesIO
from pptx.util import Inches, Pt
from django.http import HttpResponse
from django.template.loader import render_to_string
from tracker.models import FistulaCase
from mpdsr.models import MPDSREvent
from weasyprint import HTML


@login_required
def export_pdf(request):
    fistula_operated = FistulaCase.objects.filter(
        referral_status='OPERATED').count()
    total_deaths = MPDSREvent.objects.count()
    implemented = MPDSREvent.objects.filter(
        action_status='IMPLEMENTED').count()
    action_gap_percent = (implemented / total_deaths *
                          100) if total_deaths > 0 else 0

    context = {
        'fistula_operated': fistula_operated,
        'total_deaths': total_deaths,
        'action_gap_percent': action_gap_percent,
    }

    html_string = render_to_string('one_pager.html', context)
    html = HTML(string=html_string)
    pdf = html.write_pdf()

    response = HttpResponse(pdf, content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="report.pdf"'
    return response


@login_required
def export_ppt(request):
    prs = Presentation()

    # Slide 1: Title and Executive Summary
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "CIPRB M&E Report"
    subtitle.text = "Executive Summary"

    # Slide 2: Fistula Campaign Progress
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Fistula Campaign Progress"

    tf = body_shape.text_frame
    tf.text = "Surgery Outcomes Overview"

    fistula_operated = FistulaCase.objects.filter(
        referral_status='OPERATED').count()
    fistula_total = FistulaCase.objects.count()

    p = tf.add_paragraph()
    p.text = f"Total Cases: {fistula_total}"
    p.level = 1

    p = tf.add_paragraph()
    p.text = f"Operated Cases: {fistula_operated}"
    p.level = 1

    # Adding a simple table
    rows, cols = 2, 2
    left = top = Inches(3)
    width = Inches(4)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = 'Status'
    table.cell(0, 1).text = 'Count'
    table.cell(1, 0).text = 'Operated'
    table.cell(1, 1).text = str(fistula_operated)

    # Slide 3: MPDSR Action Tracking
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide3.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "MPDSR Action Tracking"

    tf = body_shape.text_frame
    tf.text = "Data-to-Action Gap Analysis"

    total_deaths = MPDSREvent.objects.count()
    implemented = MPDSREvent.objects.filter(
        action_status='IMPLEMENTED').count()
    action_gap_percent = (implemented / total_deaths *
                          100) if total_deaths > 0 else 0

    p = tf.add_paragraph()
    p.text = f"Total MPDSR Events: {total_deaths}"
    p.level = 1

    p = tf.add_paragraph()
    p.text = f"Implemented Actions: {implemented} ({action_gap_percent:.1f}%)"
    p.level = 1

    # Save to response
    response = HttpResponse(
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="report.pptx"'

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    response.write(output.read())

    return response


@login_required
def generate_newsletter(request):
    now = timezone.now()
    month = now.month
    year = now.year

    # Generate the narrative using Gemini
    narrative = generate_newsletter_narrative(month, year)

    # Save it to the model
    newsletter = MonthlyNewsletter.objects.create(
        month=month,
        year=year,
        content=narrative
    )

    return render(request, 'newsletter_result.html', {'newsletter': newsletter})
